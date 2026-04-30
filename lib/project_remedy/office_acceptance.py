"""Shared Office acceptance checks for same-format document remediation."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from project_remedy.models import FileType


@dataclass
class OfficeCheckResult:
    rule_id: str
    description: str
    status: str  # Passed / Failed / Manual Check Needed
    details: list[str] = field(default_factory=list)
    fixable: bool = False


@dataclass
class OfficeCheckReport:
    file_path: Path
    file_type: FileType
    results: list[OfficeCheckResult] = field(default_factory=list)


@dataclass
class OfficeScreenReaderIssue:
    rule_id: str
    severity: str  # error / warning
    element: str
    description: str
    suggestion: str = ""


@dataclass
class OfficeScreenReaderResult:
    file_path: Path
    file_type: FileType
    issues: list[OfficeScreenReaderIssue] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return not any(issue.severity == "error" for issue in self.issues)


@dataclass
class OfficePackageResult:
    checked: bool
    passed: bool
    error: str = ""


@dataclass
class OfficeAcceptanceResult:
    file_path: Path
    file_type: FileType
    checker_report: OfficeCheckReport
    screen_reader_result: OfficeScreenReaderResult
    package_result: OfficePackageResult

    @property
    def checker_failures(self) -> list[OfficeCheckResult]:
        return [result for result in self.checker_report.results if result.status == "Failed"]

    @property
    def screen_reader_errors(self) -> list[OfficeScreenReaderIssue]:
        return [issue for issue in self.screen_reader_result.issues if issue.severity == "error"]

    @property
    def openable(self) -> bool:
        return self.package_result.passed

    @property
    def blocking_failure_reasons(self) -> list[str]:
        if self.openable:
            return []
        if self.package_result.checked and self.package_result.error:
            return [self.package_result.error]
        return ["package validation failed"]

    @property
    def warning_entries(self) -> list[dict[str, Any]]:
        if not self.openable:
            return []
        entries: list[dict[str, Any]] = []
        for result in self.checker_failures:
            entries.append(
                {
                    "source": "checker",
                    "rule_id": result.rule_id,
                    "description": result.description,
                    "details": list(result.details),
                    "fixable": result.fixable,
                }
            )
        for issue in self.screen_reader_errors:
            entries.append(
                {
                    "source": "screen_reader",
                    "rule_id": issue.rule_id,
                    "description": issue.description,
                    "details": [issue.element] if issue.element else [],
                    "fixable": True,
                }
            )
        return entries

    @property
    def warning_reasons(self) -> list[str]:
        if not self.openable:
            return []
        reasons: list[str] = []
        if self.checker_failures:
            reasons.append(f"{len(self.checker_failures)} checker failure(s)")
        if self.screen_reader_errors:
            reasons.append(f"{len(self.screen_reader_errors)} screen reader error(s)")
        return reasons

    @property
    def retry_reasons(self) -> list[str]:
        return list(self.warning_reasons)

    @property
    def passed(self) -> bool:
        return self.openable

    def summary(self) -> str:
        if not self.openable:
            return "; ".join(self.blocking_failure_reasons)
        if self.passed and not self.warning_reasons:
            return "checker clean, screen reader clean, package valid"
        return "; ".join(self.warning_reasons)


def evaluate_office_acceptance(
    file_path: Path,
    *,
    file_type: FileType | None = None,
) -> OfficeAcceptanceResult:
    resolved_type = file_type or _infer_file_type(file_path)
    package_result = validate_office_package(file_path, resolved_type)
    if not package_result.passed:
        checker_report = OfficeCheckReport(file_path=file_path, file_type=resolved_type, results=[])
        sr_result = OfficeScreenReaderResult(file_path=file_path, file_type=resolved_type, issues=[])
        return OfficeAcceptanceResult(
            file_path=file_path,
            file_type=resolved_type,
            checker_report=checker_report,
            screen_reader_result=sr_result,
            package_result=package_result,
        )
    checker_report = run_office_checker(file_path, resolved_type)
    sr_result = run_office_screen_reader_checks(file_path, resolved_type)
    return OfficeAcceptanceResult(
        file_path=file_path,
        file_type=resolved_type,
        checker_report=checker_report,
        screen_reader_result=sr_result,
        package_result=package_result,
    )


def run_office_checker(file_path: Path, file_type: FileType) -> OfficeCheckReport:
    if file_type == FileType.DOCX:
        return _check_docx(file_path)
    if file_type == FileType.PPTX:
        return _check_pptx(file_path)
    if file_type == FileType.XLSX:
        return _check_xlsx(file_path)
    raise ValueError(f"Unsupported Office acceptance type: {file_type}")


def run_office_screen_reader_checks(file_path: Path, file_type: FileType) -> OfficeScreenReaderResult:
    if file_type == FileType.DOCX:
        return _screen_reader_docx(file_path)
    if file_type == FileType.PPTX:
        return _screen_reader_pptx(file_path)
    if file_type == FileType.XLSX:
        return _screen_reader_xlsx(file_path)
    raise ValueError(f"Unsupported Office acceptance type: {file_type}")


def validate_office_package(file_path: Path, file_type: FileType) -> OfficePackageResult:
    try:
        if file_type == FileType.DOCX:
            from docx import Document

            Document(str(file_path))
        elif file_type == FileType.PPTX:
            from pptx import Presentation

            Presentation(str(file_path))
        elif file_type == FileType.XLSX:
            from openpyxl import load_workbook

            load_workbook(str(file_path))
        else:
            raise ValueError(f"Unsupported Office acceptance type: {file_type}")
    except Exception as exc:
        return OfficePackageResult(checked=True, passed=False, error=str(exc))
    return OfficePackageResult(checked=True, passed=True)


def _check_docx(file_path: Path) -> OfficeCheckReport:
    from docx import Document

    doc = Document(str(file_path))
    results: list[OfficeCheckResult] = []

    title = (doc.core_properties.title or "").strip()
    results.append(
        OfficeCheckResult(
            rule_id="docx-title",
            description="Document title metadata is present",
            status="Passed" if title else "Failed",
            fixable=True,
        )
    )
    language = (getattr(doc.core_properties, "language", "") or "").strip()
    results.append(
        OfficeCheckResult(
            rule_id="docx-language",
            description="Document language metadata is present",
            status="Passed" if language else "Failed",
            fixable=True,
        )
    )

    headings = [para for para in doc.paragraphs if _docx_paragraph_has_heading_structure(para)]
    results.append(
        OfficeCheckResult(
            rule_id="docx-headings",
            description="Document includes heading/title styles",
            status="Passed" if headings else "Failed",
            fixable=True,
        )
    )

    missing_table_headers = 0
    for table in doc.tables:
        if not table.rows:
            continue
        tr_pr = table.rows[0]._tr.trPr
        has_header = bool(tr_pr is not None and tr_pr.find(_qn("w:tblHeader")) is not None)
        if not has_header:
            missing_table_headers += 1
    results.append(
        OfficeCheckResult(
            rule_id="docx-table-headers",
            description="Tables expose first-row header semantics",
            status="Passed" if missing_table_headers == 0 else "Failed",
            details=[f"{missing_table_headers} table(s) missing header rows"] if missing_table_headers else [],
            fixable=True,
        )
    )

    missing_alt = 0
    for inline_shape in doc.inline_shapes:
        doc_pr = inline_shape._inline.docPr
        if not ((doc_pr.get("descr") or "").strip() or (doc_pr.get("title") or "").strip()):
            missing_alt += 1
    results.append(
        OfficeCheckResult(
            rule_id="docx-alt-text",
            description="Images contain alternate text",
            status="Passed" if missing_alt == 0 else "Failed",
            details=[f"{missing_alt} image(s) missing alternate text"] if missing_alt else [],
            fixable=True,
        )
    )
    return OfficeCheckReport(file_path=file_path, file_type=FileType.DOCX, results=results)


def _check_pptx(file_path: Path) -> OfficeCheckReport:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(file_path))
    results: list[OfficeCheckResult] = []
    title = (prs.core_properties.title or "").strip()
    results.append(
        OfficeCheckResult("pptx-title", "Presentation title metadata is present", "Passed" if title else "Failed", fixable=True)
    )
    language = (getattr(prs.core_properties, "language", "") or "").strip()
    results.append(
        OfficeCheckResult("pptx-language", "Presentation language metadata is present", "Passed" if language else "Failed", fixable=True)
    )

    missing_slide_titles = 0
    missing_alt = 0
    for slide in prs.slides:
        if not _pptx_slide_title_text(slide):
            missing_slide_titles += 1
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            c_nv_pr = shape._element.nvPicPr.cNvPr
            if not ((c_nv_pr.get("descr") or "").strip() or (c_nv_pr.get("title") or "").strip()):
                missing_alt += 1

    results.append(
        OfficeCheckResult(
            "pptx-slide-titles",
            "Slides expose titles for navigation",
            "Passed" if missing_slide_titles == 0 else "Failed",
            details=[f"{missing_slide_titles} slide(s) missing titles"] if missing_slide_titles else [],
            fixable=True,
        )
    )
    results.append(
        OfficeCheckResult(
            "pptx-alt-text",
            "Pictures contain alternate text",
            "Passed" if missing_alt == 0 else "Failed",
            details=[f"{missing_alt} picture(s) missing alternate text"] if missing_alt else [],
            fixable=True,
        )
    )
    return OfficeCheckReport(file_path=file_path, file_type=FileType.PPTX, results=results)


def _check_xlsx(file_path: Path) -> OfficeCheckReport:
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path))
    results: list[OfficeCheckResult] = []
    title = (wb.properties.title or "").strip()
    results.append(
        OfficeCheckResult("xlsx-title", "Workbook title metadata is present", "Passed" if title else "Failed", fixable=True)
    )
    language = (getattr(wb.properties, "language", "") or "").strip()
    results.append(
        OfficeCheckResult("xlsx-language", "Workbook language metadata is present", "Passed" if language else "Failed", fixable=True)
    )

    missing_header_behaviors = 0
    for ws in wb.worksheets:
        has_data = ws.max_row > 1 and ws.max_column > 1
        if not has_data:
            continue
        if ws.freeze_panes != "A2" or not ws.auto_filter.ref or ws.print_title_rows is None:
            missing_header_behaviors += 1
    results.append(
        OfficeCheckResult(
            "xlsx-header-behaviors",
            "Data sheets preserve header navigation aids",
            "Passed" if missing_header_behaviors == 0 else "Failed",
            details=[f"{missing_header_behaviors} worksheet(s) missing header behaviors"] if missing_header_behaviors else [],
            fixable=True,
        )
    )
    return OfficeCheckReport(file_path=file_path, file_type=FileType.XLSX, results=results)


def _screen_reader_docx(file_path: Path) -> OfficeScreenReaderResult:
    report = run_office_checker(file_path, FileType.DOCX)
    issues = [
        OfficeScreenReaderIssue(
            rule_id=result.rule_id,
            severity="error",
            element="document",
            description=result.description,
        )
        for result in report.results
        if result.status == "Failed"
    ]
    return OfficeScreenReaderResult(file_path=file_path, file_type=FileType.DOCX, issues=issues)


def _screen_reader_pptx(file_path: Path) -> OfficeScreenReaderResult:
    report = run_office_checker(file_path, FileType.PPTX)
    issues = [
        OfficeScreenReaderIssue(
            rule_id=result.rule_id,
            severity="error",
            element="presentation",
            description=result.description,
        )
        for result in report.results
        if result.status == "Failed"
    ]
    return OfficeScreenReaderResult(file_path=file_path, file_type=FileType.PPTX, issues=issues)


def _screen_reader_xlsx(file_path: Path) -> OfficeScreenReaderResult:
    report = run_office_checker(file_path, FileType.XLSX)
    issues = [
        OfficeScreenReaderIssue(
            rule_id=result.rule_id,
            severity="error",
            element="workbook",
            description=result.description,
        )
        for result in report.results
        if result.status == "Failed"
    ]
    return OfficeScreenReaderResult(file_path=file_path, file_type=FileType.XLSX, issues=issues)


def _infer_file_type(file_path: Path) -> FileType:
    suffix = file_path.suffix.lower()
    mapping = {
        ".docx": FileType.DOCX,
        ".pptx": FileType.PPTX,
        ".xlsx": FileType.XLSX,
    }
    try:
        return mapping[suffix]
    except KeyError as exc:
        raise ValueError(f"Unsupported Office file type: {suffix}") from exc


def _pptx_slide_title_text(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is not None:
        text = (title_shape.text or "").strip()
        if text:
            return text
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if text:
            return text.splitlines()[0].strip()
    return ""


def _docx_paragraph_has_heading_structure(paragraph) -> bool:
    style_name = (getattr(getattr(paragraph, "style", None), "name", "") or "").strip().lower()
    if style_name.startswith(("title", "heading", "accessibility title", "accessibility heading")):
        return True
    return _docx_outline_level(paragraph) is not None


def _docx_outline_level(paragraph) -> int | None:
    p_pr = paragraph._p.pPr
    if p_pr is None:
        return None
    outline = p_pr.find(_qn("w:outlineLvl"))
    if outline is None:
        return None
    value = outline.get(_qn("w:val"))
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _qn(tag: str) -> str:
    prefix, local = tag.split(":")
    namespaces = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    }
    return f"{{{namespaces[prefix]}}}{local}"
